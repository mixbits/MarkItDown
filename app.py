from flask import Flask, request, render_template, send_file, flash, redirect, url_for, session
import os
import tempfile
import zipfile
from werkzeug.utils import secure_filename
import io
from datetime import datetime, timedelta
import logging
import shutil
import uuid
import time
from urllib.parse import urlparse
from flask_session import Session

# Custom MarkItDown fallback implementation
class MarkItDownResult:
    def __init__(self, text_content):
        self.text_content = text_content

class MarkItDown:
    def __init__(self, enable_plugins=False):
        self.enable_plugins = enable_plugins
        
    def convert(self, file_path):
        """Convert a file to markdown"""
        try:
            file_extension = os.path.splitext(file_path)[1].lower()
            
            if file_extension == '.txt':
                return self._convert_txt(file_path)
            elif file_extension == '.rtf':
                return self._convert_rtf(file_path)
            elif file_extension == '.pdf':
                return self._convert_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                return self._convert_docx(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return self._convert_xlsx(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return self._convert_pptx(file_path)
            elif file_extension in ['.html', '.htm']:
                return self._convert_html(file_path)
            elif file_extension == '.csv':
                return self._convert_csv(file_path)
            elif file_extension == '.json':
                return self._convert_json(file_path)
            elif file_extension == '.xml':
                return self._convert_xml(file_path)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp']:
                return self._convert_image(file_path)
            elif file_extension in ['.md', '.markdown']:
                return self._convert_markdown(file_path)
            else:
                # Fallback to text conversion
                return self._convert_txt(file_path)
                
        except Exception as e:
            error_msg = f"Error converting file: {str(e)}"
            return MarkItDownResult(error_msg)
    
    def convert_uri(self, uri):
        """Convert a URI/URL to markdown"""
        try:
            import requests
            from bs4 import BeautifulSoup
            
            # Handle YouTube URLs
            if 'youtube.com' in uri or 'youtu.be' in uri:
                return self._convert_youtube(uri)
            
            # Regular web page
            response = requests.get(uri, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and basic structure
            title = soup.find('title')
            title_text = title.get_text() if title else "Web Page"
            
            # Get main content
            content = soup.get_text()
            lines = (line.strip() for line in content.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            markdown = f"# {title_text}\n\n{text}"
            return MarkItDownResult(markdown)
            
        except Exception as e:
            error_msg = f"Error converting URL: {str(e)}"
            return MarkItDownResult(error_msg)
    
    def _convert_txt(self, file_path):
        """Convert text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return MarkItDownResult(content)
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
            return MarkItDownResult(content)
    
    def _convert_rtf(self, file_path):
        """Convert RTF using striprtf"""
        try:
            from striprtf.striprtf import rtf_to_text
            
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
            
            # Convert RTF to plain text
            text = rtf_to_text(rtf_content)
            
            # Clean up the text
            lines = text.split('\n')
            cleaned_lines = []
            for line in lines:
                line = line.strip()
                if line:  # Skip empty lines
                    cleaned_lines.append(line)
            
            cleaned_text = '\n\n'.join(cleaned_lines)
            return MarkItDownResult(cleaned_text)
            
        except ImportError:
            # Fallback: try to extract text manually from RTF
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    rtf_content = f.read()
                
                # Simple RTF text extraction (basic fallback)
                import re
                # Remove RTF control words and groups
                text = re.sub(r'\\[a-z]+\d*\s?', '', rtf_content)
                text = re.sub(r'[{}]', '', text)
                text = re.sub(r'\s+', ' ', text)
                text = text.strip()
                
                return MarkItDownResult(text)
            except Exception as fallback_error:
                return MarkItDownResult(f"Error converting RTF: striprtf library not available and fallback failed: {str(fallback_error)}. Please install striprtf: pip install striprtf")
        except Exception as e:
            return MarkItDownResult(f"Error converting RTF: {str(e)}")
    
    def _convert_pdf(self, file_path):
        """Convert PDF using pdfminer with maximum compatibility"""
        try:
            from pdfminer.high_level import extract_text
            import os
            
            # Check file size (limit to 125MB for PDF processing)
            file_size = os.path.getsize(file_path)
            if file_size > 125 * 1024 * 1024:  # 125MB limit
                return MarkItDownResult(f"Error: PDF file too large ({file_size // (1024*1024)}MB). Maximum size is 125MB.")
            
            # Try different extraction methods for maximum compatibility
            text = None
            
            # Method 1: Basic extraction with no parameters (most compatible)
            try:
                logger.info(f"Attempting basic PDF extraction for {os.path.basename(file_path)}")
                text = extract_text(file_path)
            except Exception as e1:
                logger.info(f"Basic extraction failed: {e1}")
                
                # Method 2: Try with just maxpages parameter
                try:
                    logger.info("Attempting PDF extraction with maxpages parameter")
                    text = extract_text(file_path, maxpages=100)
                except Exception as e2:
                    logger.info(f"Maxpages extraction failed: {e2}")
                    
                    # Method 3: Try with alternative low-level approach
                    try:
                        logger.info("Attempting low-level PDF extraction")
                        from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
                        from pdfminer.converter import TextConverter
                        from pdfminer.layout import LAParams
                        from pdfminer.pdfpage import PDFPage
                        from io import StringIO
                        
                        output = StringIO()
                        manager = PDFResourceManager()
                        converter = TextConverter(manager, output, laparams=LAParams())
                        interpreter = PDFPageInterpreter(manager, converter)
                        
                        with open(file_path, 'rb') as infile:
                            page_count = 0
                            for page in PDFPage.get_pages(infile, check_extractable=True):
                                interpreter.process_page(page)
                                page_count += 1
                                if page_count >= 50:  # Limit to first 50 pages
                                    break
                        
                        text = output.getvalue()
                        converter.close()
                        output.close()
                        
                    except Exception as e3:
                        logger.error(f"All PDF extraction methods failed: Basic={e1}, Maxpages={e2}, LowLevel={e3}")
                        return MarkItDownResult(f"Error: Could not extract text from PDF. All extraction methods failed. The file may be corrupted, password-protected, or contain only images.")
            
            if not text or text.strip() == '':
                return MarkItDownResult("Warning: No text could be extracted from this PDF. The PDF might contain only images or be password protected.")
            
            # Clean up the extracted text
            lines = text.split('\n')
            cleaned_lines = []
            for line in lines:
                line = line.strip()
                if line:  # Skip empty lines
                    cleaned_lines.append(line)
            
            cleaned_text = '\n\n'.join(cleaned_lines)
            logger.info(f"Successfully extracted {len(cleaned_text)} characters from PDF")
            return MarkItDownResult(cleaned_text)
            
        except ImportError:
            return MarkItDownResult("Error: PDF processing library not available. Please install pdfminer.six.")
        except Exception as e:
            logger.error(f"PDF conversion error for {file_path}: {str(e)}")
            return MarkItDownResult(f"Error converting PDF: {str(e)}. This may be due to a corrupted file, password protection, or unsupported PDF format.")
    
    def _convert_docx(self, file_path):
        """Convert DOCX using python-docx"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            markdown = ""
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Basic heading detection
                    if paragraph.style.name.startswith('Heading'):
                        level = paragraph.style.name.replace('Heading ', '')
                        if level.isdigit():
                            markdown += f"{'#' * int(level)} {text}\n\n"
                        else:
                            markdown += f"## {text}\n\n"
                    else:
                        markdown += f"{text}\n\n"
            
            # Process tables
            for table in doc.tables:
                markdown += "\n"
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    markdown += "| " + " | ".join(cells) + " |\n"
                    if i == 0:  # Header row
                        markdown += "| " + " | ".join(["---"] * len(cells)) + " |\n"
                markdown += "\n"
            
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting DOCX: {str(e)}")
    
    def _convert_xlsx(self, file_path):
        """Convert Excel using openpyxl"""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path)
            
            markdown = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                markdown += f"# {sheet_name}\n\n"
                
                # Get all rows with data
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows.append([str(cell) if cell is not None else "" for cell in row])
                
                if rows:
                    # Create table
                    for i, row in enumerate(rows[:100]):  # Limit to first 100 rows
                        markdown += "| " + " | ".join(row) + " |\n"
                        if i == 0:  # Header row
                            markdown += "| " + " | ".join(["---"] * len(row)) + " |\n"
                    markdown += "\n"
            
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting Excel: {str(e)}")
    
    def _convert_pptx(self, file_path):
        """Convert PowerPoint using python-pptx"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            markdown = "# Presentation\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                markdown += f"## Slide {i}\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            markdown += f"{text}\n\n"
            
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting PowerPoint: {str(e)}")
    
    def _convert_html(self, file_path):
        """Convert HTML using BeautifulSoup"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return MarkItDownResult(text)
        except Exception as e:
            return MarkItDownResult(f"Error converting HTML: {str(e)}")
    
    def _convert_csv(self, file_path):
        """Convert CSV using pandas"""
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            markdown = df.to_markdown(index=False)
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting CSV: {str(e)}")
    
    def _convert_json(self, file_path):
        """Convert JSON to markdown"""
        try:
            import json
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            markdown = "# JSON Data\n\n```json\n" + json.dumps(data, indent=2) + "\n```"
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting JSON: {str(e)}")
    
    def _convert_xml(self, file_path):
        """Convert XML using lxml"""
        try:
            from lxml import etree
            
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Pretty print XML
            root = etree.fromstring(content.encode())
            pretty_xml = etree.tostring(root, pretty_print=True, encoding='unicode')
            
            markdown = "# XML Data\n\n```xml\n" + pretty_xml + "\n```"
            return MarkItDownResult(markdown)
        except Exception as e:
            return MarkItDownResult(f"Error converting XML: {str(e)}")
    
    def _convert_image(self, file_path):
        """Convert image using EasyOCR (optimized for Synology NAS)"""
        try:
            import os
            import tempfile
            
            filename = os.path.basename(file_path)
            logger.info(f"Processing image: {filename}")
            
            # Set up environment for EasyOCR
            os.environ['OPENCV_IO_ENABLE_OPENEXR'] = '0'
            os.environ['DISPLAY'] = ''
            
            # EasyOCR extraction
            try:
                logger.info("Attempting OCR with EasyOCR...")
                
                # Create EasyOCR model directory
                easyocr_model_dir = os.environ.get('EASYOCR_MODULE_PATH', os.path.join(os.getcwd(), 'models', 'easyocr'))
                os.makedirs(easyocr_model_dir, exist_ok=True)
                logger.info(f"Using EasyOCR model directory: {easyocr_model_dir}")
                
                # Import and initialize EasyOCR
                import easyocr
                reader = easyocr.Reader(['en'], gpu=False, verbose=True, 
                                      model_storage_directory=easyocr_model_dir)
                
                # Try multiple parameter combinations for better results
                parameter_sets = [
                    {'detail': 1, 'paragraph': False},  # Standard
                    {'detail': 1, 'paragraph': False, 'width_ths': 0.7, 'height_ths': 0.7},  # More sensitive
                    {'detail': 1, 'paragraph': False, 'width_ths': 0.5, 'height_ths': 0.5},  # Very sensitive
                ]
                
                best_results = []
                for i, params in enumerate(parameter_sets):
                    try:
                        logger.info(f"EasyOCR attempt {i+1}/3 with params: {params}")
                        results = reader.readtext(file_path, **params)
                        logger.info(f"EasyOCR attempt {i+1} found {len(results)} text regions")
                        
                        if len(results) > len(best_results):
                            best_results = results
                            logger.info(f"New best result set with {len(results)} regions")
                        
                        if results:  # If we found something, break early
                            break
                            
                    except Exception as e:
                        logger.warning(f"EasyOCR attempt {i+1} failed: {str(e)}")
                        continue
                
                if best_results:
                    extracted_texts = []
                    logger.info(f"Processing {len(best_results)} detected text regions:")
                    
                    for i, (bbox, text, confidence) in enumerate(best_results):
                        logger.info(f"  Region {i+1}: '{text}' (confidence: {confidence:.3f})")
                        
                        if confidence > 0.1:  # Low threshold for maximum text capture
                            cleaned_text = text.strip()
                            if cleaned_text and len(cleaned_text) > 1:  # At least 2 characters
                                extracted_texts.append(cleaned_text)
                                logger.info(f"    ✅ ACCEPTED: '{cleaned_text}'")
                            else:
                                logger.info(f"    ❌ REJECTED: Too short after cleaning")
                        else:
                            logger.info(f"    ❌ REJECTED: Confidence {confidence:.3f} < 0.1")
                    
                    if extracted_texts:
                        # Join with double newlines for better formatting
                        ocr_text = '\n\n'.join(extracted_texts)
                        logger.info(f"✅ EasyOCR SUCCESS: Extracted {len(extracted_texts)} text blocks")
                        logger.info(f"   Combined text ({len(ocr_text)} chars): '{ocr_text[:200]}...'")
                        
                        # Clean up
                        del reader
                        import gc
                        gc.collect()
                        logger.info("EasyOCR cleanup completed")
                        
                        # Log final OCR summary
                        logger.info(f"🎉 OCR SUCCESS for {filename}: {len(ocr_text)} characters extracted")
                        logger.info(f"Image analysis complete for {filename}")
                        
                        # Return only the extracted text
                        return MarkItDownResult(ocr_text)
                    else:
                        logger.warning("❌ EasyOCR found text regions but none met acceptance criteria")
                        # Log all raw detections for debugging
                        for i, (bbox, text, confidence) in enumerate(best_results):
                            logger.warning(f"   Raw detection {i+1}: '{text}' (conf: {confidence:.3f})")
                else:
                    logger.warning("❌ EasyOCR found no text regions in any attempt")
                
                # Clean up on failure
                try:
                    del reader
                    import gc
                    gc.collect()
                except:
                    pass
                
            except ImportError:
                logger.warning("EasyOCR not available - install with: pip install easyocr")
            except Exception as e:
                error_msg = str(e)
                logger.error(f"EasyOCR failed with error: {error_msg}")
                import traceback
                logger.error(f"EasyOCR traceback: {traceback.format_exc()}")
                
                # Clean up on failure
                try:
                    import gc
                    gc.collect()
                except:
                    pass
            
            # If OCR failed, return error message
            logger.warning(f"❌ OCR FAILED for {filename}: No text could be extracted")
            logger.info(f"Image analysis complete for {filename}")
            return MarkItDownResult("No text could be extracted from this image.")
                
        except Exception as e:
            logger.error(f"Image conversion failed: {str(e)}")
            return MarkItDownResult(f"Error converting image: {str(e)}")
    
    def _convert_markdown(self, file_path):
        """Read existing markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return MarkItDownResult(content)
        except Exception as e:
            return MarkItDownResult(f"Error reading Markdown: {str(e)}")
    
    def _convert_youtube(self, uri):
        """Convert YouTube video (extract transcript if available)"""
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            import re
            
            # Extract video ID
            video_id = None
            if 'youtu.be' in uri:
                video_id = uri.split('/')[-1].split('?')[0]
            elif 'youtube.com' in uri:
                match = re.search(r'v=([^&]+)', uri)
                if match:
                    video_id = match.group(1)
            
            if not video_id:
                return MarkItDownResult("Error: Could not extract YouTube video ID")
            
            # Try to get transcript
            try:
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                text = ' '.join([entry['text'] for entry in transcript])
                markdown = f"# YouTube Video Transcript\n\nVideo: {uri}\n\n{text}"
                return MarkItDownResult(markdown)
            except:
                return MarkItDownResult(f"# YouTube Video\n\nVideo: {uri}\n\nTranscript not available.")
                
        except Exception as e:
            return MarkItDownResult(f"Error processing YouTube video: {str(e)}")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("logs/markitdown.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Configure session handling
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'your-secret-key-change-this-in-production')
app.config['SESSION_TYPE'] = 'filesystem'
app.config['SESSION_FILE_DIR'] = os.environ.get('SESSION_FILE_DIR', os.path.join(os.getcwd(), 'sessions'))
app.config['PERMANENT_SESSION_LIFETIME'] = timedelta(hours=1)
os.makedirs(app.config['SESSION_FILE_DIR'], exist_ok=True)
os.makedirs('logs', exist_ok=True)
Session(app)

# Configure upload settings
UPLOAD_FOLDER = os.environ.get('TMPDIR', os.path.join(os.getcwd(), 'tmp'))
TEMP_DIR = os.path.join(UPLOAD_FOLDER, 'markitdown_conversions')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)

MAX_FILE_SIZE = 250 * 1024 * 1024  # 250MB max file size

# Expanded file type support based on available libraries
ALLOWED_EXTENSIONS = {
    # Office Documents
    'pdf', 'docx', 'doc', 'pptx', 'ppt', 'xlsx', 'xls',
    # Web & Text formats  
    'html', 'htm', 'txt', 'rtf', 'csv', 'json', 'xml',
    # Image formats (OCR support)
    'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp',
    # OpenDocument formats (basic support)
    'odt', 'odp', 'ods',
    # Archive formats
    'zip',
    # Markdown (for processing)
    'md', 'markdown'
}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# Initialize custom MarkItDown converter
try:
    md_converter = MarkItDown(enable_plugins=False)
    logger.info("Custom MarkItDown converter initialized successfully")
except Exception as e:
    logger.error(f"Error initializing MarkItDown converter: {str(e)}")
    md_converter = None

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def cleanup_old_files():
    """Clean up temporary files older than 1 hour"""
    try:
        current_time = time.time()
        for dir_name in os.listdir(TEMP_DIR):
            dir_path = os.path.join(TEMP_DIR, dir_name)
            if os.path.isdir(dir_path):
                if os.path.getmtime(dir_path) < current_time - 3600:  # 1 hour
                    try:
                        shutil.rmtree(dir_path)
                        logger.info(f"Cleaned up directory {dir_path}")
                    except Exception as e:
                        logger.error(f"Error cleaning up directory {dir_path}: {str(e)}")
    except Exception as e:
        logger.error(f"Error during cleanup: {str(e)}")

def process_zip_file(zip_path, session_dir):
    """Process a ZIP file and convert all supported files within it"""
    results = {}
    
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        file_list = [f for f in zip_ref.namelist() if not f.endswith('/')]
        zip_ref.extractall(session_dir)
        
        for file_path in file_list:
            try:
                extracted_file_path = os.path.join(session_dir, file_path)
                
                if os.path.isdir(extracted_file_path) or os.path.basename(extracted_file_path).startswith('.'):
                    continue
                
                if not allowed_file(os.path.basename(file_path)):
                    logger.warning(f"Skipping unsupported file: {file_path}")
                    continue
                
                # Convert the file
                conversion_result = md_converter.convert(extracted_file_path)
                markdown_content = conversion_result.text_content
                output_filename = os.path.splitext(os.path.basename(file_path))[0] + '.md'
                
                # Save the result
                output_path = os.path.join(session_dir, output_filename)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(markdown_content)
                
                results[output_filename] = markdown_content
                logger.info(f"Successfully converted {file_path} to {output_filename}")
                
            except Exception as e:
                logger.error(f"Error processing file {file_path}: {str(e)}")
                results[os.path.basename(file_path)] = f"Error: {str(e)}"
    
    return results

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        try:
            # Handle URL conversion
            if 'url' in request.form and request.form['url'].strip():
                url = request.form['url'].strip()
                logger.info(f"Processing URL: {url}")
                
                try:
                    # Use MarkItDown's convert_uri method for URLs
                    conversion_result = md_converter.convert_uri(url)
                    result_markdown = conversion_result.text_content
                    
                    # Create session for URL result
                    session_id = str(uuid.uuid4())
                    session_dir = os.path.join(TEMP_DIR, session_id)
                    os.makedirs(session_dir, exist_ok=True)
                    
                    # Generate filename from URL
                    parsed_url = urlparse(url)
                    url_path = parsed_url.path
                    if url_path and url_path != '/':
                        base_name = os.path.basename(url_path)
                        if not base_name:
                            base_name = "url_content"
                    else:
                        base_name = parsed_url.netloc.replace('.', '_') or "url_content"
                    
                    output_filename = f"{base_name}.md"
                    output_path = os.path.join(session_dir, output_filename)
                    
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(result_markdown)
                    
                    # Store in session
                    session['conversion_id'] = session_id
                    session['single_file'] = output_filename
                    session.permanent = True
                    
                    # Return single file download
                    return send_file(
                        io.BytesIO(result_markdown.encode('utf-8')),
                        mimetype='text/markdown',
                        as_attachment=True,
                        download_name=output_filename
                    )
                    
                except Exception as e:
                    logger.error(f"Error converting URL {url}: {str(e)}")
                    flash(f'Error converting URL: {str(e)}', 'error')
                    return redirect(request.url)
            
            # Handle file uploads
            elif 'files' in request.files:
                files = request.files.getlist('files')
                if not files or all(file.filename == '' for file in files):
                    flash('No files selected', 'error')
                    return redirect(request.url)
                
                converted_files = []
                session_id = str(uuid.uuid4())
                session_dir = os.path.join(TEMP_DIR, session_id)
                os.makedirs(session_dir, exist_ok=True)
                
                for file in files:
                    if file and file.filename != '':
                        if not allowed_file(file.filename):
                            flash(f'File type not supported: {file.filename}', 'error')
                            continue
                        
                        filename = secure_filename(file.filename)
                        file_path = os.path.join(session_dir, filename)
                        file.save(file_path)
                        
                        try:
                            # Special handling for ZIP files
                            if filename.lower().endswith('.zip'):
                                zip_results = process_zip_file(file_path, session_dir)
                                for zip_filename, content in zip_results.items():
                                    converted_files.append({
                                        'original': filename,
                                        'converted': zip_filename,
                                        'path': os.path.join(session_dir, zip_filename)
                                    })
                                # Remove the ZIP file after processing
                                os.remove(file_path)
                            else:
                                # Regular file conversion
                                conversion_result = md_converter.convert(file_path)
                                result_markdown = conversion_result.text_content
                                
                                # Save markdown file
                                output_filename = os.path.splitext(filename)[0] + '.md'
                                output_path = os.path.join(session_dir, output_filename)
                                
                                with open(output_path, 'w', encoding='utf-8') as f:
                                    f.write(result_markdown)
                                
                                converted_files.append({
                                    'original': filename,
                                    'converted': output_filename,
                                    'path': output_path
                                })
                                
                                logger.info(f"Successfully converted {filename} to {output_filename}")
                        
                        except Exception as e:
                            logger.error(f"Error converting {filename}: {str(e)}")
                            flash(f'Error converting {filename}: {str(e)}', 'error')
                
                if not converted_files:
                    flash('No files could be converted', 'error')
                    return redirect(request.url)
                
                # Store conversion results in session
                session['conversion_id'] = session_id
                session['converted_files'] = [f['converted'] for f in converted_files]
                session.permanent = True
                
                # If only one file, return it directly
                if len(converted_files) == 1:
                    session['single_file'] = converted_files[0]['converted']
                    with open(converted_files[0]['path'], 'r', encoding='utf-8') as f:
                        markdown_content = f.read()
                    
                    return send_file(
                        io.BytesIO(markdown_content.encode('utf-8')),
                        mimetype='text/markdown',
                        as_attachment=True,
                        download_name=converted_files[0]['converted']
                    )
                
                # Multiple files - create ZIP
                zip_path = os.path.join(session_dir, 'converted_files.zip')
                with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for file_info in converted_files:
                        zipf.write(file_info['path'], file_info['converted'])
                
                session['zip_file'] = 'converted_files.zip'
                return send_file(
                    zip_path,
                    mimetype='application/zip',
                    as_attachment=True,
                    download_name='converted_files.zip'
                )
        
        except Exception as e:
            logger.error(f"Unexpected error in file processing: {str(e)}")
            flash(f'An unexpected error occurred: {str(e)}', 'error')
            return redirect(request.url)
    
    # GET request - show the upload form
    return render_template('index.html', 
                         allowed_extensions=sorted(ALLOWED_EXTENSIONS),
                         max_file_size_mb=MAX_FILE_SIZE // (1024*1024))

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint for monitoring"""
    try:
        # Basic health checks
        status = {
            'status': 'healthy',
            'timestamp': datetime.now().isoformat(),
            'version': '2.0.0',
            'features': {
                'file_conversion': True,
                'url_conversion': True,
                'zip_processing': True,
                'session_management': True
            },
            'supported_formats': len(ALLOWED_EXTENSIONS)
        }
        return status, 200
    except Exception as e:
        return {'status': 'unhealthy', 'error': str(e)}, 500

@app.route('/download/<filename>')
def download_file(filename):
    """Download individual converted files"""
    if 'conversion_id' not in session:
        flash('Session expired. Please convert files again.', 'error')
        return redirect(url_for('index'))
    
    session_dir = os.path.join(TEMP_DIR, session['conversion_id'])
    file_path = os.path.join(session_dir, filename)
    
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        flash('File not found or session expired', 'error')
        return redirect(url_for('index'))

@app.errorhandler(413)
def too_large(e):
    flash(f'File too large. Maximum size is {MAX_FILE_SIZE // (1024*1024)}MB', 'error')
    return redirect(request.url)

@app.before_request
def periodic_cleanup():
    """Run cleanup periodically"""
    import random
    if random.randint(1, 20) == 1:  # 5% chance
        cleanup_old_files()

@app.route('/convert_async', methods=['POST'])
def convert_async():
    """API endpoint for async conversion"""
    try:
        # Handle URL conversion
        if 'url' in request.form and request.form['url'].strip():
            url = request.form['url'].strip()
            logger.info(f"Processing URL via async: {url}")
            
            try:
                # Use MarkItDown's convert_uri method for URLs
                conversion_result = md_converter.convert_uri(url)
                result_markdown = conversion_result.text_content
                
                # Check if conversion was successful
                if result_markdown.startswith("Error"):
                    logger.error(f"URL conversion failed for {url}: {result_markdown}")
                    return {'error': result_markdown}, 400
                
                # Create session for URL result
                session_id = str(uuid.uuid4())
                session_dir = os.path.join(TEMP_DIR, session_id)
                os.makedirs(session_dir, exist_ok=True)
                
                # Generate filename from URL
                parsed_url = urlparse(url)
                url_path = parsed_url.path
                if url_path and url_path != '/':
                    base_name = os.path.basename(url_path)
                    if not base_name:
                        base_name = "url_content"
                else:
                    base_name = parsed_url.netloc.replace('.', '_') or "url_content"
                
                output_filename = f"{base_name}.md"
                output_path = os.path.join(session_dir, output_filename)
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(result_markdown)
                
                logger.info(f"URL conversion successful for {url}: {len(result_markdown)} characters extracted")
                
                # Return the file directly as a download
                return send_file(
                    io.BytesIO(result_markdown.encode('utf-8')),
                    mimetype='text/markdown',
                    as_attachment=True,
                    download_name=output_filename
                )
                
            except Exception as e:
                logger.error(f"Error converting URL {url}: {str(e)}")
                return {'error': f'Error converting URL: {str(e)}'}, 400
        
        # Handle file uploads
        # Check both 'file' and 'files' field names for compatibility
        elif 'file' in request.files:
            file = request.files['file']
        elif 'files' in request.files:
            files = request.files.getlist('files')
            if not files or files[0].filename == '':
                return {'error': 'No file provided'}, 400
            file = files[0]  # Take the first file for async processing
        else:
            return {'error': 'No file or URL provided'}, 400
            
        if file.filename == '':
            return {'error': 'No file selected'}, 400
        
        if not allowed_file(file.filename):
            return {'error': 'File type not supported'}, 400
        
        # Process file
        session_id = str(uuid.uuid4())
        session_dir = os.path.join(TEMP_DIR, session_id)
        os.makedirs(session_dir, exist_ok=True)
        
        filename = secure_filename(file.filename)
        file_path = os.path.join(session_dir, filename)
        file.save(file_path)
        
        # Convert
        logger.info(f"Starting conversion of {filename} ({os.path.getsize(file_path)} bytes)")
        conversion_result = md_converter.convert(file_path)
        result_markdown = conversion_result.text_content
        
        # Check if conversion was successful
        if result_markdown.startswith("Error") or result_markdown.startswith("Warning: No text could be extracted"):
            logger.error(f"Conversion failed for {filename}: {result_markdown}")
            return {'error': result_markdown}, 400
        
        # Log successful conversion details
        logger.info(f"Conversion successful for {filename}: {len(result_markdown)} characters extracted")
        
        # Save result and return file directly
        output_filename = os.path.splitext(filename)[0] + '.md'
        output_path = os.path.join(session_dir, output_filename)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(result_markdown)
        
        # Return the file directly as a download
        return send_file(
            io.BytesIO(result_markdown.encode('utf-8')),
            mimetype='text/markdown',
            as_attachment=True,
            download_name=output_filename
        )
    
    except Exception as e:
        logger.error(f"Error in async conversion: {str(e)}")
        return {'error': str(e)}, 500

@app.route('/download/<session_id>/<filename>')
def download_session_file(session_id, filename):
    """Download file from specific session"""
    session_dir = os.path.join(TEMP_DIR, session_id)
    file_path = os.path.join(session_dir, filename)
    
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        return {'error': 'File not found'}, 404

# Add CORS headers for API endpoints
@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
    return response

if __name__ == '__main__':
    # Get port from environment variable or default to 8008
    port = int(os.environ.get('PORT', 8008))
    
    print("=" * 60)
    print("🚀 MarkItDown Web Application")
    print("=" * 60)
    print(f"📁 Upload folder: {UPLOAD_FOLDER}")
    print(f"📁 Temp folder: {TEMP_DIR}")
    print(f"📊 Max file size: {MAX_FILE_SIZE // (1024*1024)}MB")
    print(f"📋 Supported formats: {len(ALLOWED_EXTENSIONS)} types")
    print(f"   {', '.join(sorted(ALLOWED_EXTENSIONS))}")
    print()
    print("🌐 Access URLs:")
    print(f"   Local: http://localhost:{port}")
    print(f"   Local: http://localhost:{port}")
    print(f"   Health Check: http://localhost:{port}/health")
    print(f"   Domain: https://markitdown.YOUR_DOMAIN (configure in tunnel setup)")
    print(f"   Health Check (Domain): https://health.markitdown.YOUR_DOMAIN")
    print(f"   API Access: https://api.markitdown.YOUR_DOMAIN")
    print(f"\n💡 Press Ctrl+C to stop the server\n")
    print()
    print("✨ Features:")
    print("   ✅ URL conversion support")
    print("   ✅ ZIP file processing")
    print("   ✅ Extended file format support")
    print("   ✅ Session management")
    print("   ✅ Comprehensive logging")
    print("   ✅ Health monitoring")
    print("=" * 60)
    
    app.run(host='0.0.0.0', port=port, debug=False) 